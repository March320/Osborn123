from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO

# 元数据
author_name = '熊怡嘉'
major = '财税大数据应用'
date_str = '2026-06-18'
output_filename = 'presentations/WPS演示文稿202505050214熊怡嘉.pptx'

# 示例图片（可选，若下载失败会跳过）
image_urls = [
    'https://images.unsplash.com/photo-1503676260728-1c00da094a0b?w=1200',
    'https://images.unsplash.com/photo-1521737604893-d14cc237f11d?w=1200',
    'https://images.unsplash.com/photo-1504384308090-c894fdcc538d?w=1200'
]

prs = Presentation()

# 添加封面
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = '我的大学人生规划'
if len(slide.placeholders) > 1:
    slide.placeholders[1].text = f'{author_name} • {major} • {date_str}'

# 目录
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '目录'
body = slide.shapes.placeholders[1].text_frame
body.clear()
for item in ['个人简介', '大学目标', '职业规划与能力评估', '实习与实践计划', '个人发展与技能提升', '未来展望', '致谢']:
    p = body.add_paragraph()
    p.text = item

# 若干内容页（每页控制在若干要点内）
def add_bullet(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b

add_bullet('个人简介', [
    '兴趣爱好：数据分析、阅读、旅行',
    '自我评价：积极、严谨、善于团队合作',
    '年级与背景：主修财税大数据应用'
])

add_bullet('大学目标（总览）', [
    '短期目标：提高成绩，拓展实践',
    '学业目标：通过英语/专业证书',
    '技能目标：掌握 Python/SQL/财税软件',
    '长期目标：稳定职业发展（3-5年）'
])

add_bullet('短期目标（大学期间）', [
    'GPA ≥ 3.2，完成核心课程',
    '参加比赛与社团，积累项目经验',
    '建立职业人脉，参加行业讲座'
])

add_bullet('学业目标', [
    '通过英语四/六级或同等级证书',
    '考取税务/会计/数据分析相关证书',
    '专业课程成绩持续提升'
])

add_bullet('技能提升目标', [
    '掌握 Python、SQL 与 Excel 高级技巧',
    '熟悉财税软件与数据可视化工具',
    '提升沟通与项目管理能力'
])

add_bullet('长期目标（毕业后3-5年）', [
    '进入财税或数据分析行业，担任数据分析师/税务顾问',
    '争取成为项目负责人/管理岗位',
    '实现个人生活目标（储蓄/旅行）'
])

add_bullet('职业规划与能力评估', [
    '优势：专业背景、学习能力强',
    '不足：实战经验较少、需提升业务交付',
    '行动计划：实习、项目实践、考证'
])

add_bullet('实习与实践计划', [
    '目标公司：会计/税务/金融/互联网企业',
    '时间安排：寒暑假+学期实习',
    '期望成果：完成真实数据分析或税务项目'
])

add_bullet('实习时间表（示例）', [
    '大二暑期：短期岗位体验',
    '大三学期：项目实习（数据分析）',
    '大四：毕业实习与求职准备'
])

add_bullet('个人发展计划', [
    '课程学习 + 在线课程（Coursera/慕课）',
    '项目实践：数据清洗、可视化、报表',
    '软技能提升：沟通、PPT、时间管理'
])

# 插入图片页
for i, url in enumerate(image_urls[:2]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = '项目与成果展示' if i == 0 else '致谢'
    try:
        r = requests.get(url, timeout=10)
        if r.status_code == 200:
            img = BytesIO(r.content)
            slide.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(8))
    except Exception as e:
        print('图片下载失败：', e)

add_bullet('未来展望', [
    '5 年愿景：成为资深分析师或税务顾问',
    '持续学习，保持职业敏感度与创新能力'
])

add_bullet('风险与应对', [
    '风险：岗位竞争、技术更新快',
    '对策：持续学习、建立导师与人脉'
])

add_bullet('致谢', [
    '感谢聆听',
    f'{author_name}  •  {major}  •  {date_str}'
])

# 保存演示文稿
prs.save(output_filename)
print('已生成：', output_filename)
